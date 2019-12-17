''''
Rescue metadata from embedded latexit ouptu in powerpoint slides and stick it into tex documents
'''

import glob,os
import numpy as np
import pptx
import re


def natural_sort(l): 
    convert = lambda text: int(text) if text.isdigit() else text.lower() 
    alphanum_key = lambda key: [convert(c) for c in re.split('([0-9]+)', key)] 
    return sorted(l, key=alphanum_key)

files = glob.glob('*.pptx',recursive=True)
files = [f for f in files if 'noclicker' not in f]
files = natural_sort(files)


rpre = re.compile('ESannop(.+?)ESannopend')
rsubj = re.compile('ESannot(.+?)ESannotend')

syntax_table = {'ESslash':'\\\\',
        'ESleftbrack':'{',
        'ESrightbrack':'}',
        'ESdollar':'$'}



for f in files:
    prs = pptx.Presentation(f)
    pre = []
    out = []
    for j,slide in enumerate(prs.slides):
        out.append('Slide %d\n\n'%(j+1))
        for shape in slide.shapes: 
            if isinstance(shape, pptx.shapes.picture.Picture):
                tmp = str(shape.image.blob)
                if "/Encoding /MacRomanEncoding\\n/Preamble" in tmp:

                    #update preamble
                    tmp2 = rpre.search(tmp).groups(0)[0]
                    tmp2 = tmp2.split('\\n')
                    for t in tmp2:
                        for key in syntax_table:
                            t = re.sub(key,syntax_table[key],t)
                        t = t.strip()
                        if t not in pre:
                            pre.append(t)

                    #grab subject
                    tmp2 = rsubj.search(tmp).groups(0)[0]
                    tmp2 = tmp2.split('\\n')
                    out.append('\\begin{align*}')
                    for t in tmp2:
                        for key in syntax_table:
                            t = re.sub(key,syntax_table[key],t)
                        out.append(t)
                    out.append('\\end{align*}\n\n')
                    

    with open(os.path.splitext(f)[0]+' eqs.tex', 'w') as fo:
        for p in pre:
            fo.write('%s\n' % p)
        fo.write('\\begin{document}\n\n')
        for s in out:
            fo.write(s)
        fo.write('\\end{document}\n\n')




